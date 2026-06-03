"""
PPT generator — clones slides from the reference template and fills placeholders.

Why clone instead of draw from scratch?
  The reference .pptx already has the brand fonts (Anton/Poppins), the colour
  palette, the canvas size (40 × 22.5 in), and the decorative graphics baked in.
  By cloning a template slide and only replacing the text we get a pixel-perfect
  match for free.

Theory slides are the one exception — the template has no theory layout, so we
clone the Recap layout (slide 1) and rewrite the big heading text with the
topic title. The numbered bullet boxes already match what a theory slide needs.
"""
import os
import copy
import math
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree

from schemas.slide_content import SlideContent
from schemas.slide_plan import TemplateType
from schemas.request import PDFContext
from config import OUTPUT_DIR, TEMPLATE_PPTX


# ─────────────────────────────────────────────────────────────────────────────
# TEMPLATE INDEX MAP — which slide in Common Template.pptx is the source
# (0-based; template has 14 slides total)
# ─────────────────────────────────────────────────────────────────────────────
#  0  Recap of previous lecture (orange heading + 4 numbered points + decor)
#  1  Topics to be covered      (same layout, different decorative picture)
#  2  Section heading           ("Type Heading Here" — big centred text)
#  3  Blank content slide
#  4  MCQ — vertical 4 options (A/B/C/D stacked)
#  5  MCQ — 2x2 grid of options
#  6  Question only (no options)
#  7  PYQ MCQ — vertical options, wider "Question (PYQ Exam-Year)" bar
#  8  PYQ MCQ — 2x2 grid
#  9  PYQ Question only
# 10  Summary  ("Summary" small heading + decor)
# 11  Homework ("Homework" small heading + decor)
# 12  Thank You (BLANK layout — decorative)
# 13  Style guide reference (skip)

LAYOUT_TO_TEMPLATE_IDX = {
    TemplateType.title_slide:        2,   # use section heading style for title
    TemplateType.recap_slide:        0,
    TemplateType.topics_slide:       1,
    TemplateType.section_heading:    2,
    TemplateType.theory_slide:       3,   # blank dark slide — we draw heading + bullets
    TemplateType.table_slide:        3,   # blank dark slide — we draw heading + table
    TemplateType.theory_table_slide: 3,   # blank dark slide — heading + bullets + table
    TemplateType.passage_slide:      3,   # blank dark slide — we draw banner + passage
    TemplateType.mcq_slide:          4,
    TemplateType.mcq_grid_slide:     5,
    TemplateType.question_only:      6,
    TemplateType.pyq_slide:          7,
    TemplateType.pyq_grid_slide:     8,
    TemplateType.pyq_question_only:  9,
    TemplateType.summary:           10,
    TemplateType.homework_slide:    11,
    TemplateType.thank_you_slide:   12,
}


# ─────────────────────────────────────────────────────────────────────────────
# XML helpers — clone & delete slides at the OOXML level
# ─────────────────────────────────────────────────────────────────────────────

def _clone_slide(prs, src_slide):
    """
    Deep-copy a slide's shape tree AND background image into a new slide.

    The template's slide backgrounds are embedded images (dark theme +
    decorative panels) referenced via relationship IDs (rId*). If we only
    cloned shape XML, the new slide would inherit the slide master's
    default WHITE background and our white body text would become invisible.

    So we:
      1. Copy every shape from the source slide.
      2. Copy the source slide's <p:bg> element.
      3. Walk the source slide's relationships and copy any referenced
         media parts (images) into the new slide's relationships, fixing
         up rId references inside the cloned XML.
    """
    blank_layout = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(blank_layout)

    new_cSld    = new_slide._element.find(qn('p:cSld'))
    src_cSld    = src_slide._element.find(qn('p:cSld'))
    new_sptree  = new_cSld.find(qn('p:spTree'))
    src_sptree  = src_cSld.find(qn('p:spTree'))

    # ── 1. Reset new slide's shape tree and copy shapes ─────────────────────
    for child in list(new_sptree):
        if etree.QName(child).localname not in ('nvGrpSpPr', 'grpSpPr'):
            new_sptree.remove(child)
    for child in src_sptree:
        if etree.QName(child).localname in ('nvGrpSpPr', 'grpSpPr'):
            continue
        new_sptree.append(copy.deepcopy(child))

    # ── 2. Copy <p:bg> from source slide (background image / fill) ──────────
    src_bg = src_cSld.find(qn('p:bg'))
    if src_bg is not None:
        # Remove any existing bg on the new slide first
        existing_bg = new_cSld.find(qn('p:bg'))
        if existing_bg is not None:
            new_cSld.remove(existing_bg)
        # <p:bg> must come BEFORE <p:spTree> per OOXML schema
        new_cSld.insert(list(new_cSld).index(new_sptree), copy.deepcopy(src_bg))

    # ── 3. Migrate referenced parts (background image + any embedded media) ─
    # Find every r:embed / r:link attribute inside the cloned XML and rebind
    # those rIds to fresh relationships on the new slide part.
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    src_part = src_slide.part
    new_part = new_slide.part
    rid_map: dict[str, str] = {}

    def _ensure_rid(old_rid: str) -> str | None:
        if old_rid in rid_map:
            return rid_map[old_rid]
        try:
            rel = src_part.rels[old_rid]
        except KeyError:
            return None
        # Relate the new slide part to the SAME target part as the source
        new_rid = new_part.relate_to(rel.target_part, rel.reltype)
        rid_map[old_rid] = new_rid
        return new_rid

    for elem in new_cSld.iter():
        for attr_name in (f"{{{R_NS}}}embed", f"{{{R_NS}}}link"):
            old_rid = elem.get(attr_name)
            if old_rid:
                new_rid = _ensure_rid(old_rid)
                if new_rid:
                    elem.set(attr_name, new_rid)

    return new_slide


def _delete_slides_by_indices(prs, indices):
    """Drop slides at the given 0-based indices (delete in reverse order)."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    for idx in sorted(indices, reverse=True):
        if 0 <= idx < len(slides_list):
            sl_el = slides_list[idx]
            rid = sl_el.get(qn('r:id'))
            xml_slides.remove(sl_el)
            prs.part.drop_rel(rid)


# ─────────────────────────────────────────────────────────────────────────────
# Text-replacement helpers
# ─────────────────────────────────────────────────────────────────────────────

def _iter_runs(slide):
    """Yield every text run on the slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                yield run


def _replace_first(slide, placeholder, new_value):
    """Replace the FIRST occurrence of `placeholder` with `new_value`."""
    if not new_value:
        return False
    for run in _iter_runs(slide):
        if placeholder in run.text:
            run.text = run.text.replace(placeholder, new_value)
            return True
    return False


def _replace_sequence(slide, placeholder, values):
    """
    Replace each occurrence of `placeholder` with the next item from `values`.
    Used for 4-option MCQ / numbered bullets — each is its own textbox.
    """
    it = iter(values)
    for run in _iter_runs(slide):
        if placeholder in run.text:
            try:
                run.text = run.text.replace(placeholder, next(it))
            except StopIteration:
                # leave remaining placeholders blank to avoid stray "Type option here"
                run.text = run.text.replace(placeholder, "")


def _replace_placeholders_by_shape_position(slide, placeholder, values, key):
    """
    Fill placeholder textboxes in a SPECIFIC visual order, not the XML shape order.
    `key` is a callable receiving each shape and returning a sort key.
    Used for grid layouts where XML shape order != visual order (A, C, B, D).
    """
    targets = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if placeholder in shape.text_frame.text:
            targets.append(shape)

    targets.sort(key=key)
    for shape, value in zip(targets, values):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, value)

    # blank any remaining unfilled placeholders
    for shape in targets[len(values):]:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace(placeholder, "")


def _grid_position_key(shape):
    """Sort key giving visual A, B, C, D order on a 2x2 grid."""
    # row-major: top-row first (smaller top), then left to right
    return (round(shape.top, -5), round(shape.left, -5))


def _clear_unused_placeholders(slide):
    """Blank out any leftover 'Type ... here' text so it doesn't show in output."""
    for run in _iter_runs(slide):
        if "Type option here" in run.text or "Type question here" in run.text:
            run.text = ""
        if "Type Heading Here" in run.text:
            run.text = ""


def _resolve_font_pt(shape, fallback_pt: int) -> int:
    """Return the first explicit font size on a shape, or fallback."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                return int(run.font.size.pt)
    return fallback_pt


def _apply_heading_style(
    shape,
    text_len: int,
    base_pt: int,
    min_pt: int,
    color,
    wrap: bool = True,
    max_width_in: float | None = None,
    char_width_factor: float = 0.00568,
):
    """
    Apply a heading font size that keeps the rendered text within bounds.

    When `max_width_in` is given, compute a font size such that
        text_len × pt × char_width_factor ≤ max_width_in
    so the heading text doesn't overflow its visual budget. The default
    `char_width_factor` (0.00568 in/char/pt) is calibrated against Anton at
    264pt for "Recap" (5 chars ≈ 7.5 in). Use ~0.0080 for wider fonts like
    Poppins. The template defaults (264pt big / 132pt sub) only fit the
    original "Recap" / "of previous lecture" text — any other content must
    be width-fit, not just length-banded.
    """
    from pptx.util import Pt

    if not shape:
        return
    target_pt = base_pt
    if max_width_in and text_len > 0:
        fit_pt = int(max_width_in / (text_len * char_width_factor))
        target_pt = max(min(base_pt, fit_pt), min_pt)
    elif text_len > 80:
        target_pt = max(base_pt - 28, min_pt)
    elif text_len > 65:
        target_pt = max(base_pt - 16, min_pt)
    elif text_len > 50:
        target_pt = max(base_pt - 8, min_pt)

    shape.text_frame.word_wrap = wrap
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(target_pt)
            run.font.color.rgb = color


# ─────────────────────────────────────────────────────────────────────────────
# Explicit logo cleanup
# ─────────────────────────────────────────────────────────────────────────────

_PW_BADGE_PATH = os.path.join(os.path.dirname(TEMPLATE_PPTX), "pw_badge_top_right.png")


def _remove_explicit_top_left_logo(slide) -> None:
    """
    Remove the explicit top-left PW logo cluster (white circle + PW picture).

    Body-slide templates already carry the desired branding in their cloned
    artwork/background, so adding or keeping this extra top-left cluster makes
    them diverge from the reference template. We keep the explicit logo only on
    heading-style slides; all other layouts should match the body templates.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    CIRCLE_LEFT = int(1.08 * 914400)
    CIRCLE_TOP  = int(1.00 * 914400)
    PIC_LEFT    = int(1.82 * 914400)
    PIC_TOP     = int(1.65 * 914400)
    POS_TOL     = int(0.35 * 914400)
    MAX_LOGO_W  = int(5.0  * 914400)

    for shape in list(slide.shapes):
        try:
            sw, sl, st_ = shape.width, shape.left, shape.top
        except (TypeError, AttributeError):
            continue
        if sw is None or sl is None or st_ is None or sw >= MAX_LOGO_W:
            continue

        is_circle = (
            shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and
            abs(sl - CIRCLE_LEFT) < POS_TOL and
            abs(st_ - CIRCLE_TOP) < POS_TOL
        )
        is_picture = (
            shape.shape_type == MSO_SHAPE_TYPE.PICTURE and
            abs(sl - PIC_LEFT) < POS_TOL and
            abs(st_ - PIC_TOP) < POS_TOL
        )
        if is_circle or is_picture:
            sp = shape._element
            sp.getparent().remove(sp)


def _add_top_right_badge(slide) -> None:
    """
    Add the supplied PW badge image at a fixed top-right position on the slide.

    This is the user-approved branding element to use across all generated
    slides, replacing the old top-left explicit logo cluster.
    """
    from pptx.util import Inches

    if not os.path.exists(_PW_BADGE_PATH):
        return

    # Avoid duplicate insertion if a slide gets processed twice.
    for shape in slide.shapes:
        try:
            l, t, w = shape.left, shape.top, shape.width
        except Exception:
            continue
        if abs(l - Inches(36.6)) < Inches(0.25) and abs(t - Inches(0.20)) < Inches(0.25):
            return

    slide.shapes.add_picture(
        _PW_BADGE_PATH,
        Inches(36.6),
        Inches(0.20),
        width=Inches(2.30),
    )


# ─────────────────────────────────────────────────────────────────────────────
# Per-template fillers
# ─────────────────────────────────────────────────────────────────────────────

def _strip_question_prefix(text: str) -> str:
    """Drop 'Q:', 'Question:', 'Q.1', '1.' etc. from the front of a question."""
    t = text.strip()
    for prefix in ("Question:", "Question.", "Q:", "Q.", "Ques:", "Ques.", "Problem:", "Q "):
        if t.startswith(prefix):
            t = t[len(prefix):].strip()
    # also drop leading numbering like "1." or "1) "
    if t and t[0].isdigit():
        i = 0
        while i < len(t) and (t[i].isdigit() or t[i] in ".)"):
            i += 1
        t = t[i:].strip()
    return t


def _strip_option_prefix(text: str) -> str:
    """Drop '(a)', 'a)', 'A.', etc. from each option."""
    t = text.strip()
    for pfx in (
        "(a) ", "(b) ", "(c) ", "(d) ",
        "(A) ", "(B) ", "(C) ", "(D) ",
        "a) ", "b) ", "c) ", "d) ",
        "A) ", "B) ", "C) ", "D) ",
        "a. ", "b. ", "c. ", "d. ",
        "A. ", "B. ", "C. ", "D. ",
    ):
        if t.startswith(pfx):
            return t[len(pfx):].strip()
    return t


def _fill_recap_or_topics(slide, content: SlideContent):
    """
    Slides 1/2 of template share the same shape pattern:
      - "Recap" / "Topics"      (big Anton, orange)        — first run
      - "of previous lecture" / "to be covered" (Arial)    — second run
      - 4× "Type option here"                              — bullet textboxes
    We override the heading with the slide title (split into two visual lines
    for nicer fit), and feed key_points into the bullets.
    """
    title = content.title.strip()
    # Split the title into a short "big" first chunk and a "small" remainder.
    # The big box renders at very large pt with limited horizontal room
    # (~9 in before the sub-heading box visually starts), so we cap the big
    # chunk to ~10 chars by walking word boundaries — otherwise a single long
    # first word ("Categorization") gets a tiny font and looks awkward next to
    # a fully-empty sub-heading.
    words = title.split()
    if not words:
        big, small = title, ""
    else:
        big = words[0]
        # Allow joining a second short word into the big chunk if the first
        # word is itself short, so the heading reads naturally.
        if len(words) > 1 and len(big) + 1 + len(words[1]) <= 10:
            big = f"{big} {words[1]}"
            small = " ".join(words[2:])
        else:
            small = " ".join(words[1:])

    # Find the two heading textboxes — they're the first two shapes that contain
    # the original placeholder words "Recap"/"Topics" and "of previous"/"to be".
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    heading_set = False
    sub_set = False
    heading_shape = None
    sub_shape = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        full = shape.text_frame.text.strip()
        if full in ("Recap", "Topics") and not heading_set:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.text = big
            heading_set = True
            heading_shape = shape
        elif full.startswith(("of previous", "to be")) and not sub_set:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.text = small
            sub_set = True
            sub_shape = shape

    # Width-fit the heading so longer first words don't overlap the
    # sub-heading box. Template positions: heading box L≈1.63 W≈12.96,
    # sub box L≈11.14 W≈12.96 — the boxes intentionally overlap. The big
    # text must stop visually before x≈11 (≈9 in of usable width); the sub
    # text must fit within its own ≈12 in box.
    YELLOW = RGBColor(0xFF, 0xCC, 0x31)
    base_big = _resolve_font_pt(heading_shape, 90) if heading_shape else 90
    base_small = _resolve_font_pt(sub_shape, 40) if sub_shape else 40
    _apply_heading_style(
        heading_shape,
        text_len=len(big),
        base_pt=base_big,
        min_pt=72,
        color=YELLOW,
        wrap=True,
        max_width_in=9.0,
    )
    _apply_heading_style(
        sub_shape,
        text_len=len(small),
        base_pt=base_small,
        min_pt=36,
        color=YELLOW,
        wrap=True,
        max_width_in=12.0,
        char_width_factor=0.0080,
    )

    # ── THEORY SLIDE DECOR CLEANUP ───────────────────────────────────────────
    # Theory slides reuse the Recap layout (src_idx 0). That layout has a 
    # large decorative picture (books/calculator) at the top-right.
    # We remove it for theory slides to keep them clean and professional.
    if content.layout == TemplateType.theory_slide:
        # The decorative picture in Slide 0 sits at approx (23.43, 2.37)
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        PIC_LEFT = int(23.43 * 914400)
        PIC_TOP  = int(2.37 * 914400)
        TOL      = int(1.0 * 914400)
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if abs(shape.left - PIC_LEFT) < TOL and abs(shape.top - PIC_TOP) < TOL:
                    # Remove the shape from the slide's shape tree
                    sp = shape._element
                    sp.getparent().remove(sp)

    _replace_sequence(slide, "Type option here", content.bullets[:4])
    _clear_unused_placeholders(slide)


def _fill_section_heading(slide, content: SlideContent):
    """
    Replace 'Type Heading Here' with the section title.

    The template's heading is a teal rounded-rect pill (12.99 × 1.51 in) at
    (6.11, 2.24). The default 84pt font only fits very short titles like
    "Recap" — longer titles such as "Previous Year Questions (SSC CPO - 2019)"
    wrap to 3 lines and spill below the pill. We:

      1. Pick a font size that fits the title on ONE line within the pill's
         available width (with a small inner padding so text never kisses
         the rounded ends).
      2. If the resulting size would be smaller than `_MIN_TITLE_PT`, widen
         the pill horizontally up to a hard right-edge bound (so it never
         collides with the top-right PW badge) and re-fit.
      3. Grow the pill's height to accommodate a 2-line wrap only when even
         the widened pill can't render the title at a readable size.
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    text = (content.title or "").strip()
    _replace_first(slide, "Type Heading Here", text)

    YELLOW = RGBColor(0xFF, 0xCC, 0x31)

    target_shape = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() == text:
            target_shape = shape
            break

    if target_shape is None:
        _clear_unused_placeholders(slide)
        return

    base_pt = _resolve_font_pt(target_shape, 84)
    text_len = max(len(text), 1)

    # Poppins-bold renders wider than Anton; use a font-appropriate factor.
    CHAR_W = 0.0080
    INNER_PAD_IN = 1.2        # leave room for rounded ends
    _ONE_LINE_MIN_PT = 56     # below this on one line, prefer a 2-line wrap
    _ABS_MIN_TITLE_PT = 36    # last-resort floor

    cur_width_in = target_shape.width / 914400
    cur_height_in = target_shape.height / 914400

    def _fit_pt(width_in: float, lines: int = 1) -> int:
        usable = max(width_in - INNER_PAD_IN, 1.0)
        chars_per_line = max(text_len / lines, 1)
        return int(usable / (chars_per_line * CHAR_W))

    # 1. Try one line at the pill's existing width.
    fit_pt = _fit_pt(cur_width_in, lines=1)
    use_two_lines = False

    if fit_pt < _ONE_LINE_MIN_PT:
        # 2. Two-line wrap at the same width — usually yields a larger, more
        #    readable font than cramming everything onto one tiny line.
        two_line_pt = _fit_pt(cur_width_in, lines=2)
        if two_line_pt > fit_pt:
            fit_pt = two_line_pt
            use_two_lines = True

    target_pt = max(min(base_pt, fit_pt), _ABS_MIN_TITLE_PT)

    # Grow the pill's height only when we deliberately wrap to two lines, so
    # the second line stays inside the teal rounded rectangle.
    if use_two_lines:
        # Line-height ~1.15× font size plus a small vertical pad.
        new_h_in = max(cur_height_in, target_pt / 72.0 * 2 * 1.15 + 0.30)
        target_shape.height = Inches(new_h_in)

    tf = target_shape.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(target_pt)
            run.font.color.rgb = YELLOW

    _clear_unused_placeholders(slide)


_SUBBULLET_RE = None  # lazy-built in _fill_theory_slide


def _strip_theory_prefix(text: str) -> str:
    """Drop the writer-injected '-> ' (or '➤ ') marker so the renderer owns the arrow."""
    t = text.strip()
    for pfx in ("-> ", "->", "➤ ", "➤", "• ", "•"):
        if t.startswith(pfx):
            return t[len(pfx):].lstrip()
    return t


def _clear_bullet_props(pPr):
    """Remove any existing bullet/indent child elements so we can re-set them."""
    from pptx.oxml.ns import qn
    for tag in ("a:buClr", "a:buClrTx", "a:buSzPct", "a:buSzPts", "a:buSzTx",
                "a:buFont", "a:buFontTx", "a:buChar", "a:buNone", "a:buAutoNum",
                "a:tabLst"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)


def _set_arrow_bullet(paragraph, indent_in: float, color_hex: str = "FFCC31"):
    """
    Apply a NATIVE PowerPoint bullet (➤) with a hanging indent.

    Native bullets are the reliable way to get the behaviour a hand-made deck
    has: the arrow sits at the left, the text starts at `indent_in`, and every
    WRAPPED line aligns with the text — not under the arrow. LibreOffice honours
    this where a manual "arrow run + tab" does not.

    marL = indent_in (text + wrapped lines), indent = -indent_in (bullet hangs).
    """
    from pptx.util import Inches
    from pptx.oxml.ns import qn

    pPr = paragraph._p.get_or_add_pPr()
    marL = int(Inches(indent_in))
    pPr.set("marL", str(marL))
    pPr.set("indent", str(-marL))
    _clear_bullet_props(pPr)
    # order matters (schema): buClr, buFont, buChar — appended after spcAft.
    bu_clr = etree.SubElement(pPr, qn("a:buClr"))
    etree.SubElement(bu_clr, qn("a:srgbClr")).set("val", color_hex)
    etree.SubElement(pPr, qn("a:buFont")).set("typeface", "Arial")
    etree.SubElement(pPr, qn("a:buChar")).set("char", "➤")


def _set_plain_hanging(paragraph, indent_in: float, hang_in: float):
    """Hanging indent with NO bullet glyph — for sub-bullets whose '(a)' is text."""
    from pptx.util import Inches
    from pptx.oxml.ns import qn

    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", str(int(Inches(indent_in))))
    pPr.set("indent", str(-int(Inches(hang_in))))
    _clear_bullet_props(pPr)
    etree.SubElement(pPr, qn("a:buNone"))


def _fill_theory_slide(slide, content: SlideContent):
    """
    Theory / concept layout built on top of the blank dark slide (template idx 3).

    Layout:
      - Compact yellow rounded-rect tag at top-left holding the title in black bold.
        Width auto-fits the title length.
      - Body textbox with arrow (➤) bullets in white. Bullets prefixed with
        "(a) ", "(b) ", "(c) ", "(d) " auto-indent as sub-bullets without arrow.

    The writer's "-> " prefix is stripped here so the visual arrow stays the
    renderer's responsibility.
    """
    import re
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    YELLOW = RGBColor(0xFF, 0xCC, 0x31)
    BLACK  = RGBColor(0x10, 0x10, 0x10)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

    sub_re = re.compile(r'^\(\s*([a-dA-D])\s*\)\s*')

    # ── Title — yellow tag, auto-sized to text ────────────────────────────────
    raw_title = (content.title or "").strip()
    title = raw_title.upper() if raw_title else "TOPIC"

    # Pick a title font size that fits within ~26 in (leaves room for the
    # top-right PW badge). Use a conservative char-width factor because
    # rendering engines (LibreOffice in particular) often substitute Anton
    # with a wider bold sans — under-estimating clips the text.
    MAX_TAG_W = 26.0
    PAD_X = 0.6
    PAD_Y = 0.20
    SAFETY = 0.5   # extra width padding so the title never touches the right edge
    char_w = 0.012
    for candidate_pt in (72, 64, 56, 48, 42, 36):
        if len(title) * candidate_pt * char_w <= MAX_TAG_W - 2 * PAD_X - SAFETY:
            title_pt = candidate_pt
            break
    else:
        title_pt = 36
    text_w_in = len(title) * title_pt * char_w
    tag_w = min(MAX_TAG_W, max(text_w_in + 2 * PAD_X + SAFETY, 5.0))
    tag_h = title_pt / 72.0 + 2 * PAD_Y + 0.35
    tag_l = 1.0
    tag_t = 0.8

    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(tag_l), Inches(tag_t), Inches(tag_w), Inches(tag_h),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = YELLOW
    tag.line.fill.background()
    tag.shadow.inherit = False
    tag.adjustments[0] = 0.12  # gentler rounding

    tf = tag.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(PAD_X)
    tf.margin_right = Inches(PAD_X)
    tf.margin_top = Inches(PAD_Y)
    tf.margin_bottom = Inches(PAD_Y)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_pt)
    run.font.bold = True
    run.font.name = "Anton"
    run.font.color.rgb = BLACK

    # ── Body — arrow bullets + optional (a)/(b) sub-bullets ──────────────────
    body_left = Inches(1.5)
    body_top = Inches(tag_t + tag_h + 0.9)
    body_width = Inches(37.0)
    body_height = Inches(22.5 - (tag_t + tag_h + 0.9) - 1.2)  # leave footer room

    bullets = [b for b in (content.bullets or []) if b and b.strip()]
    if not bullets:
        return

    # Body font size comes from the shared fit engine so the generator and the
    # reflow/pagination engine always agree on capacity. The reflow pass has
    # already split this slide so its bullets fit; pick_body_font then renders
    # them at the largest font that fits the box height.
    from pipeline.fit_engine import pick_body_font
    body_pt = pick_body_font(bullets, TemplateType.theory_slide)

    body_tb = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
    bt = body_tb.text_frame
    bt.word_wrap = True
    bt.vertical_anchor = MSO_ANCHOR.TOP

    first = True
    for raw in bullets:
        text = _strip_theory_prefix(raw)
        is_sub = bool(sub_re.match(text))

        p = bt.paragraphs[0] if first else bt.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(14 if is_sub else 24)   # set BEFORE bullet/indent props

        # Hanging-indent width scales with the font. The canvas is 40in wide, so
        # the indent must be sizeable for the arrow→text gap and the wrapped-line
        # alignment to read clearly (a small 0.5in indent is invisible here).
        main_indent = round(body_pt * 0.020, 2)      # ≈ 0.9in @44pt … 1.1in @56pt

        if is_sub:
            # Sub-bullet "(a) …": indented deeper, no arrow glyph; its own hanging
            # indent so a wrapped sub-line aligns under the sub-bullet text.
            _set_plain_hanging(p, main_indent + 0.7, 0.55)
            run_t = p.add_run()
            run_t.text = text
            run_t.font.size = Pt(max(body_pt - 4, 22))
            run_t.font.name = "Poppins"
            run_t.font.color.rgb = WHITE
        else:
            # Native ➤ bullet (yellow) + hanging indent → wrapped lines align.
            _set_arrow_bullet(p, main_indent, color_hex="FFCC31")
            run_t = p.add_run()
            run_t.text = text
            run_t.font.size = Pt(body_pt)
            run_t.font.name = "Poppins"
            run_t.font.color.rgb = WHITE


# ─────────────────────────────────────────────────────────────────────────────
# Table renderers — table_slide and theory_table_slide
# ─────────────────────────────────────────────────────────────────────────────

def _draw_yellow_title_tag(slide, raw_title: str, top_in: float = 0.8) -> tuple[float, float]:
    """
    Shared helper: draws the same yellow rounded-rect title tag used by
    theory_slide at the top of the slide. Returns the (left, top + height)
    in inches, so callers know where the body content can start.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    YELLOW = RGBColor(0xFF, 0xCC, 0x31)
    BLACK  = RGBColor(0x10, 0x10, 0x10)

    title = (raw_title or "TOPIC").strip().upper() or "TOPIC"

    MAX_TAG_W = 26.0
    PAD_X = 0.6
    PAD_Y = 0.20
    SAFETY = 0.5
    char_w = 0.012
    for candidate_pt in (72, 64, 56, 48, 42, 36):
        if len(title) * candidate_pt * char_w <= MAX_TAG_W - 2 * PAD_X - SAFETY:
            title_pt = candidate_pt
            break
    else:
        title_pt = 36
    text_w_in = len(title) * title_pt * char_w
    tag_w = min(MAX_TAG_W, max(text_w_in + 2 * PAD_X + SAFETY, 5.0))
    tag_h = title_pt / 72.0 + 2 * PAD_Y + 0.35
    tag_l = 1.0

    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(tag_l), Inches(top_in), Inches(tag_w), Inches(tag_h),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = YELLOW
    tag.line.fill.background()
    tag.shadow.inherit = False
    tag.adjustments[0] = 0.12

    tf = tag.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(PAD_X)
    tf.margin_right = Inches(PAD_X)
    tf.margin_top = Inches(PAD_Y)
    tf.margin_bottom = Inches(PAD_Y)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_pt)
    run.font.bold = True
    run.font.name = "Anton"
    run.font.color.rgb = BLACK

    return tag_l, top_in + tag_h


def _looks_numeric(value: str) -> bool:
    """Heuristic: 'is this cell a number?' — used for default column alignment."""
    if value is None:
        return False
    s = str(value).strip()
    if not s:
        return False
    # strip common decorations: currency, %, parentheses, commas, leading +/-
    s = s.replace(",", "").replace("$", "").replace("₹", "").replace("%", "")
    s = s.replace("(", "").replace(")", "").strip()
    if not s:
        return False
    try:
        float(s)
        return True
    except ValueError:
        return False


def _pick_table_font_size(rows_count: int, cols_count: int,
                          available_height_in: float,
                          longest_cell_len: int,
                          available_col_width_in: float) -> int:
    """
    Pick a font size that lets `rows_count` rows fit in `available_height_in`
    AND lets the longest cell fit within `available_col_width_in`.

    Returns a pt size in [12, 32].
    """
    # Height-based cap: each row is roughly font_pt * 1.6 (incl. inner padding)
    # in points → inches = pt / 72.
    if rows_count <= 0:
        return 22
    h_pt = (available_height_in * 72.0) / (rows_count * 1.6)

    # Width-based cap: longest cell text shouldn't overflow its column.
    char_w_factor = 0.0080  # Poppins approx, inches per (char × pt)
    if longest_cell_len > 0 and available_col_width_in > 0:
        w_pt = available_col_width_in / (longest_cell_len * char_w_factor)
    else:
        w_pt = 999

    candidate = int(min(h_pt, w_pt))
    return max(12, min(candidate, 32))


def _add_styled_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    column_alignments: list[str] | None = None,
):
    """
    Render a real PowerPoint table on the dark slide.

    Visual language:
      - Header row : solid yellow (#FFCC31) fill, black bold text.
      - Body rows  : alternating very-dark fills (#1F1F1F / #2A2A2A) so each
                     row reads cleanly against the dark template background.
                     White text.
      - Borders    : thin dark grey lines so the grid is visible but not loud.
      - Column widths: proportional to the longest cell length in each column.
      - Font size  : auto-shrunk so all rows fit in `height_in`.

    `column_alignments` is an optional list with "left"/"center"/"right" per
    column. When omitted, numeric-looking columns default to right-aligned
    and text columns default to left-aligned.
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    YELLOW   = RGBColor(0xFF, 0xCC, 0x31)
    BLACK    = RGBColor(0x10, 0x10, 0x10)
    WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
    ROW_DARK = RGBColor(0x1F, 0x1F, 0x1F)
    ROW_MID  = RGBColor(0x2A, 0x2A, 0x2A)
    BORDER   = RGBColor(0x55, 0x55, 0x55)

    cols = len(headers)
    if cols == 0 or not rows:
        return
    # Defensive: pad/trim every row to len(headers).
    norm_rows = []
    for r in rows:
        r = list(r) + [""] * max(0, cols - len(r))
        norm_rows.append([str(c) if c is not None else "" for c in r[:cols]])

    total_rows = len(norm_rows) + 1  # +1 for header

    table_shape = slide.shapes.add_table(
        total_rows, cols,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    table = table_shape.table

    # ── Column widths from content -----------------------------------------
    col_text_len = [max(len(headers[c]), 1) for c in range(cols)]
    for r in norm_rows:
        for c in range(cols):
            col_text_len[c] = max(col_text_len[c], len(r[c]))
    total_len = sum(col_text_len) or 1
    # Reserve a small minimum width (1.0 in) per column so single-character cells
    # don't collapse to nothing.
    MIN_COL_IN = 1.0
    base_widths = [max(MIN_COL_IN, width_in * (l / total_len)) for l in col_text_len]
    # Rescale to actually equal width_in.
    scale = width_in / sum(base_widths)
    widths = [w * scale for w in base_widths]
    for c in range(cols):
        table.columns[c].width = Inches(widths[c])

    # ── Font size auto-fit -------------------------------------------------
    longest_cell = max(col_text_len)
    avg_col_w = sum(widths) / cols
    font_pt = _pick_table_font_size(
        rows_count=total_rows,
        cols_count=cols,
        available_height_in=height_in,
        longest_cell_len=longest_cell,
        available_col_width_in=avg_col_w,
    )

    # Row heights — distribute the total height roughly evenly. The header
    # gets a 1.15× share so it reads as a distinct band.
    header_share = 1.15
    unit = height_in / (header_share + (total_rows - 1))
    table.rows[0].height = Inches(unit * header_share)
    for r in range(1, total_rows):
        table.rows[r].height = Inches(unit)

    # ── Default alignment per column ---------------------------------------
    if column_alignments and len(column_alignments) == cols:
        col_align = column_alignments
    else:
        col_align = []
        for c in range(cols):
            # numeric if >= 60% of body cells in this column parse as numbers
            n = sum(1 for r in norm_rows if _looks_numeric(r[c]))
            col_align.append("right" if (n / max(len(norm_rows), 1)) >= 0.6 else "left")

    def _apply_align(p, name):
        if name == "right":
            p.alignment = PP_ALIGN.RIGHT
        elif name == "center":
            p.alignment = PP_ALIGN.CENTER
        else:
            p.alignment = PP_ALIGN.LEFT

    # ── Header row ---------------------------------------------------------
    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = YELLOW
        cell.margin_left = Inches(0.10)
        cell.margin_right = Inches(0.10)
        cell.margin_top = Inches(0.06)
        cell.margin_bottom = Inches(0.06)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf = cell.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        _apply_align(p, "center")  # headers always centered
        run = p.add_run()
        run.text = headers[c] or ""
        run.font.size = Pt(font_pt)
        run.font.bold = True
        run.font.name = "Poppins"
        run.font.color.rgb = BLACK

    # ── Body rows ----------------------------------------------------------
    for r_idx, row in enumerate(norm_rows):
        bg = ROW_DARK if r_idx % 2 == 0 else ROW_MID
        for c in range(cols):
            cell = table.cell(r_idx + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.margin_left = Inches(0.10)
            cell.margin_right = Inches(0.10)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            tf = cell.text_frame
            tf.word_wrap = True
            tf.clear()
            p = tf.paragraphs[0]
            _apply_align(p, col_align[c])
            run = p.add_run()
            run.text = row[c] or ""
            run.font.size = Pt(max(font_pt - 1, 12))
            run.font.name = "Poppins"
            run.font.color.rgb = WHITE

    return table_shape


def _draw_table_caption(slide, caption: str, left_in: float, top_in: float,
                        width_in: float) -> float:
    """Optional small italic caption above the table. Returns its bottom Y."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if not caption:
        return top_in
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.6),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = caption.strip()
    run.font.size = Pt(24)
    run.font.italic = True
    run.font.name = "Poppins"
    run.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    return top_in + 0.55


def _fill_table_slide(slide, content: SlideContent):
    """
    Table-only slide on the blank dark template (idx 3).

    Layout:
      - Yellow rounded-rect title tag at top-left (same style as theory_slide)
      - Optional italic caption immediately below the tag
      - A real PowerPoint table filling the rest of the body area

    Falls back to `_fill_theory_slide` (using the slide's bullets) if the
    writer didn't supply structured table_data — that way the slide is never
    empty even when the writer fails.
    """
    table_data = getattr(content, "table_data", None)
    if not table_data or not table_data.headers or not table_data.rows:
        # No usable table — let theory slide handle it as a graceful degrade.
        _fill_theory_slide(slide, content)
        return

    _, body_top = _draw_yellow_title_tag(slide, content.title, top_in=0.8)

    LEFT = 1.0
    WIDTH = 38.0           # canvas is 40 in → leaves 1in margin each side
    BOTTOM_LIMIT = 21.4    # leave room for the context footer at the bottom

    cur_top = body_top + 0.45
    cur_top = _draw_table_caption(
        slide, table_data.caption or "", LEFT, cur_top, WIDTH,
    )

    table_h = max(BOTTOM_LIMIT - cur_top, 3.0)
    _add_styled_table(
        slide,
        headers=list(table_data.headers),
        rows=[list(r) for r in table_data.rows],
        left_in=LEFT,
        top_in=cur_top,
        width_in=WIDTH,
        height_in=table_h,
        column_alignments=table_data.column_alignments,
    )


def _fill_theory_table_slide(slide, content: SlideContent):
    """
    Theory bullets ABOVE a small table on the blank dark template (idx 3).

    Layout (top → bottom):
      1. Yellow title tag
      2. Theory bullets (arrow ➤ style, same as theory_slide)
      3. Optional table caption
      4. Real PowerPoint table

    The renderer guarantees no overlap: bullet block height is bounded so the
    table always has at least ~6 inches of vertical room. If the bullets would
    push the table off the slide, the body font shrinks first; if the table is
    still too tall, the renderer downgrades to a table-only slide (drops the
    bullets) so the data — which is harder to compress — stays readable.
    """
    import re
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    table_data = getattr(content, "table_data", None)
    bullets = [b for b in (content.bullets or []) if b and b.strip()]

    # If there's no table at all, fall through to the plain theory layout.
    if not table_data or not table_data.headers or not table_data.rows:
        _fill_theory_slide(slide, content)
        return

    # If there are no bullets, prefer the cleaner table-only layout.
    if not bullets:
        _fill_table_slide(slide, content)
        return

    _, body_top = _draw_yellow_title_tag(slide, content.title, top_in=0.8)

    BODY_LEFT  = 1.5
    BODY_WIDTH = 37.0
    BOTTOM_LIMIT = 21.4

    # Bullets block — give the bullets a bounded chunk of vertical space so
    # they never crowd the table. With at most 3 bullets at 32pt (after fit),
    # ~3.5in is plenty.
    bullets_top = body_top + 0.5
    BULLETS_MAX_H = 5.0    # in inches
    MIN_TABLE_H  = 6.0     # the table always gets at least this much room

    available_after_bullets = BOTTOM_LIMIT - (bullets_top + BULLETS_MAX_H + 0.6)
    if available_after_bullets < MIN_TABLE_H:
        # Compress the bullets allowance.
        BULLETS_MAX_H = max(2.5, BOTTOM_LIMIT - bullets_top - MIN_TABLE_H - 0.6)

    sub_re = re.compile(r'^\(\s*([a-dA-D])\s*\)\s*')

    # Body font for the bullet block — modest so 2-3 bullets fit in BULLETS_MAX_H.
    bullet_pt = 28 if len(bullets) >= 3 else 32

    body_tb = slide.shapes.add_textbox(
        Inches(BODY_LEFT), Inches(bullets_top),
        Inches(BODY_WIDTH), Inches(BULLETS_MAX_H),
    )
    bt = body_tb.text_frame
    bt.word_wrap = True
    bt.vertical_anchor = MSO_ANCHOR.TOP

    first = True
    for raw in bullets[:3]:
        text = _strip_theory_prefix(raw)
        is_sub = bool(sub_re.match(text))
        p = bt.paragraphs[0] if first else bt.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(10 if is_sub else 18)
        main_indent = round(bullet_pt * 0.020, 2)
        if is_sub:
            _set_plain_hanging(p, main_indent + 0.7, 0.55)
            run_t = p.add_run()
            run_t.text = text
            run_t.font.size = Pt(max(bullet_pt - 4, 22))
            run_t.font.name = "Poppins"
            run_t.font.color.rgb = WHITE
        else:
            _set_arrow_bullet(p, main_indent, color_hex="FFCC31")
            run_t = p.add_run()
            run_t.text = text
            run_t.font.size = Pt(bullet_pt)
            run_t.font.name = "Poppins"
            run_t.font.color.rgb = WHITE

    # Table sits below the bullets block.
    table_top = bullets_top + BULLETS_MAX_H + 0.4
    table_top = _draw_table_caption(
        slide, table_data.caption or "", BODY_LEFT, table_top, BODY_WIDTH,
    )
    table_h = max(BOTTOM_LIMIT - table_top, MIN_TABLE_H)

    _add_styled_table(
        slide,
        headers=list(table_data.headers),
        rows=[list(r) for r in table_data.rows],
        left_in=BODY_LEFT,
        top_in=table_top,
        width_in=BODY_WIDTH,
        height_in=table_h,
        column_alignments=table_data.column_alignments,
    )


def _fill_passage_slide(slide, content: SlideContent):
    """
    Cloze / reading-comprehension passage on the blank dark slide (template idx 3).

    Visual language matches the theory slide (same color theme):
      - A yellow rounded-rect banner at top holding the `directions` line in
        black bold (e.g. "Directions (Q. 22-24): Cloze Test – Passage 1").
      - The verbatim passage paragraph below in white, justified, with every
        blank (__X__, .....(1).....) preserved EXACTLY as written by the writer.

    The passage is rendered as ONE flowing, atomic paragraph — never split or
    paraphrased — so the gaps stay visible for the student to fill.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pipeline.fit_engine import estimate_block_height_in

    YELLOW = RGBColor(0xFF, 0xCC, 0x31)
    BLACK  = RGBColor(0x10, 0x10, 0x10)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

    # ── Directions banner — yellow rounded rect, black bold, word-wrapped ─────
    directions = (content.directions or content.title or "Passage").strip()

    band_l, band_t = 1.0, 0.8
    band_w = 36.0                       # wide banner; leaves room for PW badge
    band_pad_x, band_pad_y = 0.5, 0.18

    # Fit the directions font to the banner width (word-wrap to 1-2 lines).
    char_w = 0.0095                     # in/char/pt for the bold sans banner
    usable_w = band_w - 2 * band_pad_x
    dir_pt = 40
    for candidate_pt in (40, 36, 32, 28, 24):
        chars_per_line = max(1, int(usable_w / (candidate_pt * char_w)))
        if len(directions) <= chars_per_line * 2:   # fits within two lines
            dir_pt = candidate_pt
            break
    else:
        dir_pt = 24
    chars_per_line = max(1, int(usable_w / (dir_pt * char_w)))
    dir_lines = max(1, math.ceil(len(directions) / chars_per_line))
    band_h = (dir_pt / 72.0) * dir_lines * 1.25 + 2 * band_pad_y

    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(band_l), Inches(band_t), Inches(band_w), Inches(band_h),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = YELLOW
    band.line.fill.background()
    band.shadow.inherit = False
    band.adjustments[0] = 0.18

    btf = band.text_frame
    btf.word_wrap = True
    btf.margin_left = Inches(band_pad_x)
    btf.margin_right = Inches(band_pad_x)
    btf.margin_top = Inches(band_pad_y)
    btf.margin_bottom = Inches(band_pad_y)
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.LEFT
    brun = bp.add_run()
    brun.text = directions
    brun.font.size = Pt(dir_pt)
    brun.font.bold = True
    brun.font.name = "Anton"
    brun.font.color.rgb = BLACK

    # ── Passage body — white, justified, verbatim (blanks preserved) ──────────
    passage = (content.passage_text or "").strip()
    if not passage:
        # Fallback: writer didn't populate passage_text — join any bullets so
        # nothing is lost (still verbatim, just not pre-formatted).
        passage = "\n".join(b for b in (content.bullets or []) if b and b.strip())
    if not passage:
        return

    body_left = Inches(1.5)
    body_top = Inches(band_t + band_h + 0.7)
    body_width_in = 37.0
    body_top_in = band_t + band_h + 0.7
    body_height_in = max(6.0, 22.5 - body_top_in - 1.2)   # leave footer room

    # FILL THE SLIDE like a hand-made deck — not small text clustered at the top:
    #   1) pick the LARGEST font (big canvas ⇒ up to ~56pt) that still fits;
    #   2) spread the lines (line spacing) to use the leftover height;
    #   3) vertically CENTRE so any residual gap is balanced, never all at bottom.
    paras = [ln for ln in passage.split("\n") if ln.strip()] or [passage]
    pass_pt = 30
    for pt in range(56, 29, -2):
        if estimate_block_height_in(paras, pt, body_width_in) <= body_height_in:
            pass_pt = pt
            break
    else:
        pass_pt = 30

    # Spread lines to fill the height. line_spacing ≤ 1.3 × (box / natural) keeps
    # the block from overflowing while expanding it toward the full height.
    natural = estimate_block_height_in(paras, pass_pt, body_width_in)
    line_spacing = 1.3
    if natural > 0:
        line_spacing = max(1.3, min(1.55, 1.3 * (body_height_in / natural)))

    body_tb = slide.shapes.add_textbox(
        body_left, body_top, Inches(body_width_in), Inches(body_height_in)
    )
    bt = body_tb.text_frame
    bt.word_wrap = True
    bt.vertical_anchor = MSO_ANCHOR.MIDDLE   # centre vertically → fills, no top-cluster

    for i, para_text in enumerate(paras):
        p = bt.paragraphs[0] if i == 0 else bt.add_paragraph()
        p.alignment = PP_ALIGN.JUSTIFY
        p.line_spacing = line_spacing
        p.space_after = Pt(18)
        run = p.add_run()
        run.text = para_text
        run.font.size = Pt(pass_pt)
        run.font.name = "Poppins"
        run.font.color.rgb = WHITE


def _add_context_footer(slide, context: PDFContext):
    """
    Small context strip bottom-left on body slides — keeps subject / batch /
    purpose visible across the deck without crowding the template's layout.
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    parts = [p for p in (context.subject, context.batch, context.purpose) if p]
    if not parts:
        return
    text = "  ·  ".join(parts)

    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(21.6), Inches(30.0), Inches(0.7)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(20)
    run.font.italic = True
    run.font.name = "Poppins"
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def _fill_title_slide(slide, content: SlideContent, context: PDFContext):
    """
    Build a proper title slide on top of the section-heading layout.
    The template only has one heading textbox, so we add three more textboxes
    around it to show the full lecture context:
      - main heading  : the lecture / topic title  (replaces 'Type Heading Here')
      - subtitle      : Subject  ·  Batch
      - metadata      : Purpose  ·  Class Level  ·  Language
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    YELLOW = RGBColor(0xFF, 0xCC, 0x31)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY   = RGBColor(0xAA, 0xAA, 0xAA)

    # 1. Main heading — the lecture title
    topic = content.title or f"{context.subject} — {context.purpose}"
    _replace_first(slide, "Type Heading Here", topic)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() == topic:
            base_pt = _resolve_font_pt(shape, 84)
            _apply_heading_style(
                shape,
                text_len=len(topic),
                base_pt=base_pt,
                min_pt=56,
                color=YELLOW,
                wrap=True,
            )
            break

    # 2. Subtitle — Subject · Batch
    subtitle_parts = []
    if context.subject:
        subtitle_parts.append(context.subject)
    if context.batch:
        subtitle_parts.append(context.batch)
    subtitle = "  ·  ".join(subtitle_parts)

    if subtitle:
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(10.0), Inches(38.0), Inches(2.0)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(60)
        run.font.bold = True
        run.font.name = "Poppins"
        run.font.color.rgb = WHITE

    # 3. Metadata — Purpose · Class Level · Language
    meta_parts = []
    if context.purpose:
        meta_parts.append(context.purpose)
    if context.class_level:
        meta_parts.append(context.class_level)
    if context.language:
        meta_parts.append(context.language)
    metadata = "    ·    ".join(meta_parts)

    if metadata:
        tb = slide.shapes.add_textbox(
            Inches(1.0), Inches(13.5), Inches(38.0), Inches(1.2)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = metadata
        run.font.size = Pt(36)
        run.font.name = "Poppins"
        run.font.color.rgb = GRAY

    # 4. Bottom accent bar — yellow line
    from pptx.util import Emu
    bar = slide.shapes.add_shape(
        1,  # rectangle
        Inches(8.0), Inches(20.5), Inches(24.0), Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = YELLOW
    bar.line.fill.background()

    _clear_unused_placeholders(slide)


def _sanitize_question_title(text: str) -> str:
    """
    Safety net: strip any 'Answer:', 'Exam:', exam year tags, or explanation
    text that the writer may have accidentally put in the title field.
    """
    import re
    t = text.strip()
    # Cut at "Answer:" if present — everything after is the answer
    t = re.split(r'\s*Answer\s*:', t, maxsplit=1, flags=re.IGNORECASE)[0].strip()
    # Cut at "Exam:" if present — that's exam metadata
    t = re.split(r'\s*Exam\s*:', t, maxsplit=1, flags=re.IGNORECASE)[0].strip()
    # Cut at standalone "(SSC" or "(Exam" pattern — year tag in parens
    t = re.split(r'\s*\(\s*(?:SSC|Exam|PYQ|JEE|NEET|UPSC)', t, maxsplit=1)[0].strip()
    return t


def _fill_mcq(slide, content: SlideContent, is_grid: bool = False):
    q = _strip_question_prefix(content.title)
    q = _sanitize_question_title(q)
    _replace_first(slide, "Type question here", q)
    opts = [_strip_option_prefix(b) for b in content.bullets[:4]]

    if is_grid:
        # In the template the 4 option boxes are arranged COL1-row1, COL1-row2,
        # COL2-row1, COL2-row2 — i.e. A, C, B, D in XML order. We need to fill
        # them in row-major visual order so the user sees A, B, C, D correctly.
        _replace_placeholders_by_shape_position(
            slide, "Type option here", opts, key=_grid_position_key
        )
    else:
        _replace_sequence(slide, "Type option here", opts)

    _clear_unused_placeholders(slide)


def _extract_exam_tag(notes: str) -> str | None:
    """
    Pull ONLY the exam name/year from speaker_notes, stripping any answer
    text that may have leaked onto the same line.

    Expected format:  "Exam: SSC CGL Tier-II 11/09/2019\nAnswer: (a) ..."
    But we also handle:  "Exam: SSC CGL 2019 Answer: (a) Dirge" (no newline).
    """
    if not notes:
        return None
    for line in notes.splitlines():
        line = line.strip()
        if not line.lower().startswith(("exam", "pyq")):
            continue
        # Chop off anything starting with "Answer" on the same line
        import re
        tag = re.split(r'\s*Answer\s*:', line, maxsplit=1, flags=re.IGNORECASE)[0].strip()
        # Remove the "Exam:" prefix itself to keep just the name + year
        for pfx in ("Exam:", "Exam :", "PYQ:", "PYQ :", "Exam-", "Exam"):
            if tag.startswith(pfx):
                tag = tag[len(pfx):].strip()
                break
        if tag:
            # Truncate to keep the banner from overflowing
            if len(tag) > 50:
                tag = tag[:47] + "..."
            return tag
    return None


def _fill_pyq(slide, content: SlideContent, is_grid: bool = False):
    """Same as MCQ but also fills the PYQ subtitle if speaker_notes carries it."""
    _fill_mcq(slide, content, is_grid=is_grid)
    tag = _extract_exam_tag(content.speaker_notes)
    if tag:
        banner_text = f"Question (Exam: {tag})"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full = shape.text_frame.text
            if "Question (Type PYQ" in full or "Question (PYQ" in full:
                # set the banner text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "Question (Type PYQ" in run.text or "Question (PYQ" in run.text:
                            run.text = banner_text
                # auto-shrink font so the banner never overflows its box
                # default template font is ~32-36pt; reduce for long tags
                total_len = len(banner_text)
                if   total_len > 65: target_pt = 20
                elif total_len > 52: target_pt = 24
                elif total_len > 40: target_pt = 28
                else:                target_pt = None   # keep template default
                if target_pt:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(target_pt)
                # also enable word-wrap so nothing spills outside the box
                shape.text_frame.word_wrap = True
                break


def _fill_question_only(slide, content: SlideContent):
    from pptx.util import Inches
    from pptx.dml.color import RGBColor

    q = _strip_question_prefix(content.title)
    q = _sanitize_question_title(q)
    _replace_first(slide, "Type question here", q)
    _clear_unused_placeholders(slide)

    # If the writer provided bullets/options (happens when layout was set to
    # question_only but the question actually has MCQ options), render them as
    # a 2-column option grid below the question area so nothing is lost.
    opts = [_strip_option_prefix(b) for b in content.bullets if b.strip()]
    if not opts:
        return

    labels = ["(A)", "(B)", "(C)", "(D)"]
    col_w    = Inches(17.5)
    col_gap  = Inches(3.0)
    left_x   = [Inches(1.5), Inches(1.5) + col_w + col_gap]
    row_h    = Inches(2.8)
    start_y  = Inches(10.5)   # below the question text area

    for i, opt in enumerate(opts[:4]):
        col = i % 2
        row = i // 2
        tb = slide.shapes.add_textbox(
            left_x[col],
            start_y + row * row_h,
            col_w,
            row_h,
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{labels[i]}  {opt}"
        run.font.size = Pt(34)
        run.font.name = "Poppins"
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _fill_thank_you(slide, content: SlideContent):
    # The thank-you slide is fully decorative — nothing to replace.
    _clear_unused_placeholders(slide)


def _add_bullets_textbox(slide, bullets, top_in=6.0, font_pt=40):
    """
    Append a simple bullets textbox below the heading. Used for summary /
    homework where the template only has a small title and no body area.
    Coordinates are in the template's 40 × 22.5 in canvas.
    `font_pt` is supplied by the fit engine so the size matches the layout's
    capacity (the reflow pass has already paginated overflow).
    """
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    tb = slide.shapes.add_textbox(
        Inches(1.5), Inches(top_in), Inches(37.0), Inches(15.0)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(18)
        run = p.add_run()
        run.text = f"{i + 1}.  {line}"
        run.font.size = Pt(font_pt)
        run.font.name = "Poppins"
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _fill_summary_or_homework(slide, content: SlideContent):
    """Template only has heading; add a body textbox for the points."""
    if content.bullets:
        from pipeline.fit_engine import pick_body_font
        font_pt = pick_body_font(content.bullets, content.layout)
        _add_bullets_textbox(slide, content.bullets, top_in=6.0, font_pt=font_pt)


# ─────────────────────────────────────────────────────────────────────────────
# Speaker notes
# ─────────────────────────────────────────────────────────────────────────────

def _set_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


# ─────────────────────────────────────────────────────────────────────────────
# Router — pick the right filler for each layout
# ─────────────────────────────────────────────────────────────────────────────

def _apply_content(slide, content: SlideContent, context: PDFContext):
    t = content.layout
    if t == TemplateType.title_slide:
        _fill_title_slide(slide, content, context)
    elif t == TemplateType.recap_slide:
        _fill_recap_or_topics(slide, content)
    elif t == TemplateType.topics_slide:
        _fill_recap_or_topics(slide, content)
    elif t == TemplateType.section_heading:
        _fill_section_heading(slide, content)
    elif t == TemplateType.theory_slide:
        # theory uses the new yellow-tag + arrow-bullets layout on a blank base
        _fill_theory_slide(slide, content)
    elif t == TemplateType.table_slide:
        # table-only — yellow caption tag at top, real pptx table fills the body
        _fill_table_slide(slide, content)
    elif t == TemplateType.theory_table_slide:
        # short theory bullets above + table below, non-overlapping
        _fill_theory_table_slide(slide, content)
    elif t == TemplateType.passage_slide:
        # cloze/comprehension passage — yellow directions banner + verbatim text
        _fill_passage_slide(slide, content)
    elif t == TemplateType.mcq_slide:
        _fill_mcq(slide, content, is_grid=False)
    elif t == TemplateType.mcq_grid_slide:
        _fill_mcq(slide, content, is_grid=True)
    elif t == TemplateType.pyq_slide:
        _fill_pyq(slide, content, is_grid=False)
    elif t == TemplateType.pyq_grid_slide:
        _fill_pyq(slide, content, is_grid=True)
    elif t == TemplateType.question_only:
        _fill_question_only(slide, content)
    elif t == TemplateType.pyq_question_only:
        _fill_question_only(slide, content)
        # patch only the PYQ exam-tag banner — don't re-fill question/options
        tag = _extract_exam_tag(content.speaker_notes)
        if tag:
            banner_text = f"Question (Exam: {tag})"
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                full = shape.text_frame.text
                if "Question (Type PYQ" in full or "Question (PYQ" in full:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if "Question (Type PYQ" in run.text or "Question (PYQ" in run.text:
                                run.text = banner_text
                    total_len = len(banner_text)
                    if   total_len > 65: target_pt = 20
                    elif total_len > 52: target_pt = 24
                    elif total_len > 40: target_pt = 28
                    else:                target_pt = None
                    if target_pt:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.size = Pt(target_pt)
                    shape.text_frame.word_wrap = True
                    break
    elif t == TemplateType.summary:
        _fill_first_text(slide, "Summary", content.title or "Summary")
        _fill_summary_or_homework(slide, content)
    elif t == TemplateType.homework_slide:
        _fill_first_text(slide, "Homework", content.title or "Homework")
        _fill_summary_or_homework(slide, content)
    elif t == TemplateType.thank_you_slide:
        _fill_thank_you(slide, content)
    else:
        _clear_unused_placeholders(slide)

    # Keep subject/batch/purpose visible on every body slide.
    skip_footer_on = {TemplateType.title_slide, TemplateType.thank_you_slide}
    if t not in skip_footer_on:
        _add_context_footer(slide, context)

    _set_notes(slide, content.speaker_notes)


def _fill_first_text(slide, old, new):
    """Replace the first textbox whose text equals `old`."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip() == old:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.text = new
            return True
    return False


# ─────────────────────────────────────────────────────────────────────────────
# MAIN ENTRY
# ─────────────────────────────────────────────────────────────────────────────

def generate_pptx(
    all_slide_contents: list[SlideContent],
    context: PDFContext,
    filename: str = "output.pptx"
) -> str:
    """
    Build the final deck by cloning slides from the reference template and
    filling placeholders with our generated content.
    """
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    if not os.path.exists(TEMPLATE_PPTX):
        raise FileNotFoundError(
            f"Reference template not found: {TEMPLATE_PPTX}. "
            "Add 'Common Template.pptx' to assets/reference_ppts/."
        )

    prs = Presentation(TEMPLATE_PPTX)
    original_count = len(prs.slides)

    for content in all_slide_contents:
        src_idx = LAYOUT_TO_TEMPLATE_IDX.get(content.layout)
        if src_idx is None:
            # unknown template type — fall back to theory/recap layout
            src_idx = LAYOUT_TO_TEMPLATE_IDX[TemplateType.theory_slide]

        try:
            new_slide = _clone_slide(prs, prs.slides[src_idx])
            _apply_content(new_slide, content, context)
            _remove_explicit_top_left_logo(new_slide)
            _add_top_right_badge(new_slide)

            print(f"    Slide {content.slide_number:2d} [{content.layout.value:18s}] — "
                  f"{content.title[:55]}")
        except Exception as e:
            print(f"    Slide {content.slide_number:2d} — failed: {e}")

    # Drop the original 14 template slides; keep only the ones we built.
    _delete_slides_by_indices(prs, list(range(original_count)))

    out = os.path.join(OUTPUT_DIR, filename)
    prs.save(out)
    print(f"\n  PPT saved → {out}")
    return out
